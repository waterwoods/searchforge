#!/usr/bin/env python3
"""
Generate FIQA summary presentation from PDF report.
Creates a screenshot from the PDF and builds a PowerPoint presentation.
"""

import os
from pathlib import Path

try:
    import fitz  # PyMuPDF
except ImportError:
    print("Installing PyMuPDF...")
    os.system("pip install PyMuPDF -q")
    import fitz

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Installing python-pptx...")
    os.system("pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt


def extract_pdf_as_image(pdf_path: str, output_path: str, dpi: int = 300):
    """Extract PDF page as high-resolution PNG image."""
    print(f"📄 Opening PDF: {pdf_path}")
    doc = fitz.open(pdf_path)
    
    # Get first page (assuming single page report)
    page = doc[0]
    
    # Render at high DPI for quality
    zoom = dpi / 72  # 72 is the default DPI
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    
    # Save as PNG
    pix.save(output_path)
    print(f"✅ Screenshot saved: {output_path}")
    print(f"   Resolution: {pix.width}x{pix.height} pixels")
    
    doc.close()
    return output_path


def create_presentation(screenshot_path: str, output_ppt: str):
    """Create PowerPoint presentation with title and screenshot."""
    print(f"\n🎨 Creating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title page
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "Finance QA AutoTuner Canary Report"
    subtitle1.text = "FiQA 数据集金丝雀验证结果"
    
    # Style title
    title1.text_frame.paragraphs[0].font.size = Pt(44)
    title1.text_frame.paragraphs[0].font.bold = True
    
    # Slide 2: Summary card with screenshot
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "AutoTuner 一页报告结论卡片"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    
    # Add screenshot (centered)
    img_left = Inches(0.8)
    img_top = Inches(1.3)
    img_width = Inches(8.4)
    slide2.shapes.add_picture(screenshot_path, img_left, img_top, width=img_width)
    
    # Add summary text at bottom
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "系统在金丝雀阶段展现出稳定的智能调优效果 (p < 0.05, Recall ↑5%, 延迟稳定)。"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.italic = True
    
    # Save presentation
    prs.save(output_ppt)
    print(f"✅ PowerPoint created: {output_ppt}")


def main():
    # Setup paths
    project_root = Path(__file__).parent.parent
    pdf_path = project_root / "docs" / "one_pager_fiqa.pdf"
    screenshot_dir = project_root / "docs" / "screenshots"
    screenshot_path = screenshot_dir / "fiqa_summary_card.png"
    presentation_dir = project_root / "docs" / "presentations"
    ppt_path = presentation_dir / "fiqa_canary_summary.pptx"
    
    # Create directories if needed
    screenshot_dir.mkdir(exist_ok=True, parents=True)
    presentation_dir.mkdir(exist_ok=True, parents=True)
    
    # Step 1: Extract PDF as high-res image
    extract_pdf_as_image(str(pdf_path), str(screenshot_path), dpi=300)
    
    # Step 2: Create PowerPoint
    create_presentation(str(screenshot_path), str(ppt_path))
    
    print("\n" + "="*60)
    print("✨ Task completed successfully!")
    print("="*60)
    print(f"📸 Screenshot: {screenshot_path}")
    print(f"📊 Presentation: {ppt_path}")
    print("="*60)


if __name__ == "__main__":
    main()



